"""
Artifact creation and reading tools as nanobot Tool ABC subclasses.

  - CreateArtifactTool  (create_artifact)
  - ReadArtifactTool    (read_artifact)
"""

from __future__ import annotations

import json
import os
from typing import Any

from nanobot.agent.tools.base import Tool
from nanobot.clawwork.state import ClawWorkState


class CreateArtifactTool(Tool):
    """Create a work artifact file in the sandbox directory."""

    def __init__(self, state: ClawWorkState) -> None:
        self._state = state

    @property
    def name(self) -> str:
        return "create_artifact"

    @property
    def description(self) -> str:
        return (
            "Create a work artifact file (txt, md, csv, json, xlsx, docx, pdf) "
            "in the sandbox directory."
        )

    @property
    def parameters(self) -> dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "filename": {
                    "type": "string",
                    "description": "Name for the file (without extension).",
                },
                "content": {
                    "type": "string",
                    "description": (
                        "Content to write. For xlsx, provide JSON array or CSV text. "
                        "For json, provide valid JSON string."
                    ),
                },
                "file_type": {
                    "type": "string",
                    "enum": ["txt", "md", "csv", "json", "xlsx", "docx", "pdf"],
                    "description": "File format (default: txt).",
                },
            },
            "required": ["filename", "content"],
        }

    async def execute(self, **kwargs: Any) -> str:
        import json as json_lib

        filename: str = kwargs.get("filename", "")
        content: str = kwargs.get("content", "")
        file_type: str = kwargs.get("file_type", "txt").lower().strip()

        if not filename:
            return json.dumps({"error": "Filename cannot be empty"})
        if not content:
            return json.dumps({"error": "Content cannot be empty"})

        valid_types = ["txt", "md", "csv", "json", "xlsx", "docx", "pdf"]
        if file_type not in valid_types:
            return json.dumps(
                {
                    "error": f"Invalid file type: {file_type}",
                    "valid_types": valid_types,
                }
            )

        data_path = self._state.data_path
        date = self._state.current_date

        if not data_path:
            return json.dumps({"error": "Data path not configured"})

        sandbox_dir = os.path.join(data_path, "sandbox", date or "default")
        os.makedirs(sandbox_dir, exist_ok=True)

        safe_filename = os.path.basename(filename).replace("/", "_").replace("\\", "_")
        file_path = os.path.join(sandbox_dir, f"{safe_filename}.{file_type}")

        try:
            if file_type in ("txt", "md", "csv"):
                with open(file_path, "w", encoding="utf-8") as f:
                    f.write(content)

            elif file_type == "json":
                try:
                    json_data = json_lib.loads(content)
                    with open(file_path, "w", encoding="utf-8") as f:
                        json_lib.dump(json_data, f, indent=2, ensure_ascii=False)
                except json_lib.JSONDecodeError as e:
                    return json.dumps({"error": f"Invalid JSON content: {e}"})

            elif file_type == "xlsx":
                try:
                    import pandas as pd

                    try:
                        data = json_lib.loads(content)
                        df = pd.DataFrame(data)
                    except Exception:
                        import io

                        df = pd.read_csv(io.StringIO(content))
                    df.to_excel(file_path, index=False, engine="openpyxl")
                except ImportError:
                    return json.dumps(
                        {"error": "openpyxl not installed. Run: pip install openpyxl pandas"}
                    )
                except Exception as e:
                    return json.dumps({"error": f"Failed to create Excel file: {e}"})

            elif file_type == "docx":
                try:
                    from docx import Document

                    doc = Document()
                    for para in content.split("\n\n"):
                        if para.strip():
                            doc.add_paragraph(para.strip())
                    doc.save(file_path)
                except ImportError:
                    return json.dumps(
                        {"error": "python-docx not installed. Run: pip install python-docx"}
                    )
                except Exception as e:
                    return json.dumps({"error": f"Failed to create Word document: {e}"})

            elif file_type == "pdf":
                try:
                    from reportlab.lib.pagesizes import letter
                    from reportlab.lib.styles import getSampleStyleSheet
                    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

                    doc = SimpleDocTemplate(file_path, pagesize=letter)
                    styles = getSampleStyleSheet()
                    story = []
                    for para in content.split("\n\n"):
                        if para.strip():
                            story.append(Paragraph(para.strip(), styles["Normal"]))
                            story.append(Spacer(1, 12))
                    doc.build(story)
                except ImportError:
                    return json.dumps(
                        {"error": "reportlab not installed. Run: pip install reportlab"}
                    )
                except Exception as e:
                    return json.dumps({"error": f"Failed to create PDF: {e}"})

            file_size = os.path.getsize(file_path)
            return json.dumps(
                {
                    "success": True,
                    "filename": f"{safe_filename}.{file_type}",
                    "file_path": file_path,
                    "file_type": file_type,
                    "file_size": file_size,
                    "message": (
                        f"Created {file_type.upper()} file: {safe_filename}.{file_type} "
                        f"({file_size} bytes). To submit as work artifact, call "
                        f'submit_work(artifact_file_paths=["{file_path}"])'
                    ),
                }
            )

        except Exception as e:
            return json.dumps({"error": f"Failed to create file: {e}", "filename": safe_filename})


class ReadArtifactTool(Tool):
    """Read a file and return its content."""

    def __init__(self, state: ClawWorkState) -> None:
        self._state = state

    @property
    def name(self) -> str:
        return "read_artifact"

    @property
    def description(self) -> str:
        return (
            "Read a file and return its content. "
            "Supports pdf, docx, xlsx, pptx, png, jpg, jpeg, txt."
        )

    @property
    def parameters(self) -> dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "filetype": {
                    "type": "string",
                    "enum": ["pdf", "docx", "xlsx", "pptx", "png", "jpg", "jpeg", "txt"],
                    "description": "The type of file to read.",
                },
                "file_path": {
                    "type": "string",
                    "description": "Absolute path to the file.",
                },
            },
            "required": ["filetype", "file_path"],
        }

    async def execute(self, **kwargs: Any) -> str:
        from pathlib import Path

        filetype: str = kwargs.get("filetype", "").lower().strip()
        file_path_str: str = kwargs.get("file_path", "")

        if not filetype or not file_path_str:
            return json.dumps({"error": "Both filetype and file_path are required"})

        file_path = Path(file_path_str)

        if not file_path.exists():
            return json.dumps({"error": f"File not found: {file_path}"})

        supported = ("pdf", "docx", "xlsx", "pptx", "png", "jpg", "jpeg", "txt")
        if filetype not in supported:
            return json.dumps(
                {
                    "error": f"Unsupported file type: {filetype}",
                    "supported_types": list(supported),
                }
            )

        try:
            if filetype == "txt":
                with open(file_path, "r", encoding="utf-8") as f:
                    text = f.read()
                return json.dumps({"type": "text", "text": text})

            elif filetype == "docx":
                try:
                    from docx import Document

                    doc = Document(file_path)
                    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                    text = "\n\n".join(paragraphs)
                    return json.dumps({"type": "text", "text": text})
                except ImportError:
                    return json.dumps(
                        {"error": "python-docx not installed. Run: pip install python-docx"}
                    )

            elif filetype == "xlsx":
                try:
                    import pandas as pd

                    df = pd.read_excel(file_path, engine="openpyxl")
                    return json.dumps({"type": "text", "text": df.to_string()})
                except ImportError:
                    return json.dumps(
                        {"error": "openpyxl/pandas not installed. Run: pip install openpyxl pandas"}
                    )

            elif filetype in ("png", "jpg", "jpeg"):
                import base64

                with open(file_path, "rb") as f:
                    image_data = base64.b64encode(f.read()).decode("utf-8")
                return json.dumps(
                    {
                        "type": "image",
                        "format": filetype,
                        "data": image_data,
                        "message": "Image file loaded as base64.",
                    }
                )

            elif filetype == "pdf":
                try:
                    import fitz  # PyMuPDF

                    doc = fitz.open(file_path)
                    text = ""
                    for page in doc:
                        text += page.get_text()
                    return json.dumps({"type": "text", "text": text})
                except ImportError:
                    return json.dumps(
                        {"error": "PyMuPDF not installed. Run: pip install pymupdf"}
                    )

            elif filetype == "pptx":
                try:
                    from pptx import Presentation

                    prs = Presentation(file_path)
                    text_parts = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text_parts.append(shape.text)
                    return json.dumps({"type": "text", "text": "\n\n".join(text_parts)})
                except ImportError:
                    return json.dumps(
                        {"error": "python-pptx not installed. Run: pip install python-pptx"}
                    )

            else:
                return json.dumps({"error": f"Unsupported file type: {filetype}"})

        except Exception as e:
            return json.dumps({"error": f"Failed to read file: {e}"})
